#!/usr/bin/env python3
"""
Script para gerar apresentação PowerPoint: "Como a LIA Revoluciona o Recrutamento Moderno"
Baseado na documentação institucional da plataforma LIA
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def criar_apresentacao_lia():
    # Criar nova apresentação
    prs = Presentation()

    # Configurar slide master
    slide_width = Inches(13.33)  # 16:9 format
    slide_height = Inches(7.5)

    # Cores da marca LIA
    cor_primaria = RGBColor(139, 92, 246)  # Purple
    cor_secundaria = RGBColor(59, 130, 246)  # Blue
    cor_sucesso = RGBColor(16, 185, 129)  # Green
    cor_alerta = RGBColor(245, 158, 11)  # Orange
    cor_texto = RGBColor(31, 41, 55)  # Dark gray

    # Lista de slides com conteúdo
    slides_content = [
        {
            "titulo": "🚀 COMO A LIA REVOLUCIONA\nO RECRUTAMENTO MODERNO",
            "subtitulo": "Transformação Completa do RH através de IA\nDo Manual ao Inteligente | Do Reativo ao Preditivo",
            "tipo": "titulo"
        },
        {
            "titulo": "📋 COMO VAMOS REVOLUCIONAR O RH",
            "conteudo": [
                "🚨 O Problema do Recrutamento Atual",
                "💎 A Revolução da LIA",
                "⚡ Transformação em 4 Dimensões",
                "🎯 Tecnologias Revolucionárias",
                "📊 Resultados Comprovados",
                "🔄 Processo Transformado",
                "🚀 O Futuro Chegou"
            ],
            "tipo": "lista"
        },
        {
            "titulo": "🚨 A REALIDADE DOLOROSA DO RH ATUAL",
            "conteudo": [
                "📈 NÚMEROS ALARMANTES:",
                "• 95% dos CVs nunca são lidos",
                "• 60+ dias para fechar uma vaga",
                "• R$ 15.000 custo de contratação errada",
                "• 40% tempo gasto em tarefas repetitivas",
                "",
                "😰 CONSEQUÊNCIAS:",
                "• Recrutadores sobrecarregados",
                "• Decisões baseadas em 'feeling'",
                "• Perda de talentos para concorrência"
            ],
            "tipo": "problema"
        },
        {
            "titulo": "💎 LIA: A REVOLUÇÃO CHEGOU",
            "conteudo": [
                "🤖 PRIMEIRA IA CONVERSACIONAL",
                "Especializada 100% em Talent Acquisition",
                "",
                "⚡ AUTOMAÇÃO INTELIGENTE",
                "70% das tarefas executadas automaticamente",
                "",
                "🎯 PRECISÃO CIENTÍFICA",
                "95% accuracy no match candidato-vaga",
                "",
                "📊 DECISÕES DATA-DRIVEN",
                "Fim do 'feeling' - início da ciência"
            ],
            "tipo": "solucao"
        },
        {
            "titulo": "🚀 REVOLUÇÃO EM 4 DIMENSÕES",
            "conteudo": [
                "1. 🧠 INTELIGÊNCIA AUMENTADA",
                "   Antes: 'Acho que ele tem o perfil'",
                "   Agora: 'Score 94/100 - fit garantido'",
                "",
                "2. ⚡ AUTOMAÇÃO TOTAL",
                "   Antes: 8h triando 50 CVs",
                "   Agora: 15min + ações sugeridas",
                "",
                "3. 🔮 PREDIÇÃO vs. REAÇÃO",
                "   Antes: Problemas pós-contratação",
                "   Agora: Antecipação de riscos/sucessos",
                "",
                "4. 🌍 OPERAÇÃO 24/7",
                "   Antes: Processo para fora do expediente",
                "   Agora: LIA trabalhando sempre"
            ],
            "tipo": "transformacao"
        },
        {
            "titulo": "🎯 SCORE LIA: PRECISÃO CIENTÍFICA",
            "conteudo": [
                "📊 ANÁLISE MULTICRITÉRIO:",
                "• TÉCNICO: Skills + experiência + resultados",
                "• CULTURAL: Values + fit + comunicação",
                "• COMPORTAMENTAL: Big Five automatizado",
                "• PREDITIVO: Probabilidade de sucesso",
                "",
                "🧠 INTELIGÊNCIA AVANÇADA:",
                "• 1000+ candidatos processados simultaneamente",
                "• Score 0-100% com justificativas detalhadas",
                "• Identificação automática de red flags",
                "",
                "⚡ RESULTADO:",
                "95% precisão vs. 60% do mercado tradicional"
            ],
            "tipo": "feature"
        },
        {
            "titulo": "🤖 LIA CONVERSA COMO ESPECIALISTA",
            "conteudo": [
                "💬 INTERAÇÃO NATURAL:",
                "'LIA, o que você sugere para este candidato?'",
                "",
                "🎯 RESPOSTAS INTELIGENTES:",
                "'Score 89/100. Recomendo:",
                "1. Entrevista técnica hoje 15h",
                "2. Case fintech personalizado",
                "3. Alerta - está sendo headhunted'",
                "",
                "⚡ EXECUÇÃO AUTOMÁTICA:",
                "• Agendamentos coordenados",
                "• Mensagens personalizadas",
                "• Follow-ups inteligentes",
                "• Relatórios automáticos"
            ],
            "tipo": "demo"
        },
        {
            "titulo": "⚡ 19 AÇÕES AUTOMATIZADAS",
            "conteudo": [
                "🎯 TRIAGEM & ANÁLISE:",
                "• Fazer triagem completa",
                "• Coletar dados faltantes",
                "• Atualizar avaliação LIA",
                "",
                "📱 COMUNICAÇÃO:",
                "• Enviar WhatsApp personalizado",
                "• Email templates inteligentes",
                "• Agendar entrevistas",
                "",
                "📊 GESTÃO & RELATÓRIOS:",
                "• Compartilhar perfis",
                "• Gerar relatórios executivos",
                "• Transferir candidatos",
                "",
                "= 70% DO TRABALHO EXECUTADO PELA LIA"
            ],
            "tipo": "automacao"
        },
        {
            "titulo": "📊 RESULTADOS REVOLUCIONÁRIOS",
            "conteudo": [
                "⚡ EFICIÊNCIA OPERACIONAL:",
                "• 70% redução tempo triagem",
                "• 50% redução time-to-hire",
                "• 40% economia custos recrutamento",
                "• 3x produtividade recrutadores",
                "",
                "🎯 QUALIDADE SUPERIOR:",
                "• 95% precisão match candidato-vaga",
                "• 60% redução contratações inadequadas",
                "• 90% satisfação dos gestores",
                "",
                "💰 ROI COMPROVADO:",
                "• 245% ROI médio primeiro ano",
                "• Payback em 1.6 meses",
                "• R$ 800K+ economia anual típica"
            ],
            "tipo": "resultados"
        },
        {
            "titulo": "📊 CASE: MULTINACIONAL TECH",
            "conteudo": [
                "🎯 ANTES DA LIA:",
                "• 180 vagas/ano | R$ 18K cost-per-hire",
                "• 75 dias time-to-hire",
                "• 35% contratações inadequadas",
                "",
                "⚡ APÓS 12 MESES COM LIA:",
                "• Cost-per-hire: R$ 18K → R$ 9.5K (47% ↓)",
                "• Time-to-hire: 75 → 28 dias (63% ↓)",
                "• Quality: 65% → 94% (45% ↑)",
                "",
                "💰 ROI FINAL:",
                "• Economia: R$ 2.4M/ano",
                "• ROI: 1233% primeiro ano",
                "",
                "💬 'LIA transformou nosso RH de cost center em profit center'"
            ],
            "tipo": "case"
        },
        {
            "titulo": "🏆 POR QUE LIA NÃO TEM CONCORRENTE",
            "conteudo": [
                "🥇 IA CONVERSACIONAL ESPECIALIZADA",
                "❌ Outros: Chatbots genéricos",
                "✅ LIA: Especialista em RH brasileiro",
                "",
                "🥇 SCORE 360° PROPRIETÁRIO",
                "❌ Outros: Match por palavras-chave",
                "✅ LIA: Análise científica multicritério",
                "",
                "🥇 BIG FIVE AUTOMÁTICO",
                "❌ Outros: Testes manuais separados",
                "✅ LIA: Psicologia integrada ao workflow",
                "",
                "🥇 SETUP EM 48H",
                "❌ Outros: 3-6 meses implementação",
                "✅ LIA: Funcionando em 2 dias"
            ],
            "tipo": "diferencial"
        },
        {
            "titulo": "🚀 SETUP EM 48H, NÃO 6 MESES",
            "conteudo": [
                "📅 DIA 1 - SETUP TÉCNICO:",
                "• Infrastructure deployment (2h)",
                "• Data migration + integration (6h)",
                "• Security configuration (4h)",
                "",
                "📅 DIA 2 - GO-LIVE:",
                "• Team training intensivo (4h)",
                "• Process customization (2h)",
                "• First candidates processing (2h)",
                "",
                "📅 SEMANA 1 - MASTERY:",
                "• Advanced features activation",
                "• Workflow optimization",
                "• Performance monitoring",
                "",
                "⚡ RESULTADO:",
                "ROI positivo em 30 dias, não 6 meses"
            ],
            "tipo": "implementacao"
        },
        {
            "titulo": "💰 RETORNO TRANSFORMADOR",
            "conteudo": [
                "📊 ECONOMIA TÍPICA (500 funcionários):",
                "• Redução cost-per-hire: R$ 400K/ano",
                "• Aumento produtividade: R$ 300K/ano",
                "• Redução retrabalho: R$ 250K/ano",
                "",
                "💎 INVESTIMENTO LIA:",
                "• Plano Professional: R$ 90K/ano",
                "• Setup + training inclusos",
                "• Updates automáticos gratuitos",
                "",
                "🎯 RESULTADO FINAL:",
                "• Economia total: R$ 1.15M/ano",
                "• ROI líquido: R$ 1.06M/ano",
                "• Percentual: 1178% primeiro ano",
                "• Payback: 1.2 meses"
            ],
            "tipo": "roi"
        },
        {
            "titulo": "⚠️ JANELA DE OPORTUNIDADE FECHANDO",
            "conteudo": [
                "🔥 WAR FOR TALENT INTENSIFICANDO:",
                "• Talentos 70% mais escassos",
                "• Decisões em 48-72h máximo",
                "• Processos lentos = talentos perdidos",
                "",
                "📈 REVOLUÇÃO IA EM CURSO:",
                "• Early adopters dominando mercado",
                "• 2+ anos vantagem competitiva",
                "• Late adopters ficando obsoletos",
                "",
                "🚨 RISCO DE INAÇÃO:",
                "• Perder melhores talentos",
                "• Custos crescendo exponencialmente",
                "• Concorrentes tomando liderança",
                "",
                "= LIDERAR REVOLUÇÃO OU SER ULTRAPASSADO"
            ],
            "tipo": "urgencia"
        },
        {
            "titulo": "🚀 REVOLUCIONE SEU RH HOJE",
            "conteudo": [
                "🎯 DECISÃO SIMPLES:",
                "",
                "❌ CONTINUAR NO PASSADO:",
                "• Processos manuais lentos",
                "• Custos crescendo infinitamente",
                "• Perdendo war for talent",
                "",
                "✅ LIDERAR O FUTURO:",
                "• IA revolucionária exclusiva",
                "• ROI 245% comprovado",
                "• Vantagem 2+ anos sobre mercado",
                "",
                "📞 PRÓXIMO PASSO:",
                "'Quando agendamos a demo?'",
                "",
                "⚡ CONTATO: +55 11 9999-8888",
                "📧 lia@wedo.com.br"
            ],
            "tipo": "cta"
        }
    ]

    # Criar slides
    for i, slide_data in enumerate(slides_content):
        if slide_data["tipo"] == "titulo":
            criar_slide_titulo(prs, slide_data, cor_primaria, cor_texto)
        else:
            criar_slide_conteudo(prs, slide_data, cor_primaria, cor_secundaria, cor_texto)

    # Salvar apresentação
    nome_arquivo = "LIA_Revoluciona_Recrutamento_Moderno.pptx"
    prs.save(nome_arquivo)
    print(f"✅ Apresentação criada com sucesso: {nome_arquivo}")
    return nome_arquivo

def criar_slide_titulo(prs, slide_data, cor_primaria, cor_texto):
    """Criar slide de título"""
    slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(slide_layout)

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = slide_data["titulo"]
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = cor_primaria
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = slide_data["subtitulo"]
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = cor_texto
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Logo/marca (placeholder)
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "WeDo Talent © 2025 | Apresentação Institucional"
    footer_frame.paragraphs[0].font.size = Pt(14)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def criar_slide_conteudo(prs, slide_data, cor_primaria, cor_secundaria, cor_texto):
    """Criar slide de conteúdo"""
    slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(slide_layout)

    # Título do slide
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_data["titulo"]
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = cor_primaria

    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.auto_size = True

    # Adicionar conteúdo linha por linha
    for i, linha in enumerate(slide_data["conteudo"]):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = linha

        # Estilizar baseado no tipo de conteúdo
        if linha.startswith(("📊", "⚡", "🎯", "💰", "🔥", "📈")):
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = cor_secundaria
        elif linha.startswith("•"):
            p.font.size = Pt(16)
            p.font.color.rgb = cor_texto
            p.level = 1
        elif linha.startswith(("❌", "✅")):
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(220, 53, 69) if linha.startswith("❌") else RGBColor(25, 135, 84)
            p.level = 1
        elif linha == "":
            p.font.size = Pt(8)  # Espaçamento
        else:
            p.font.size = Pt(16)
            p.font.color.rgb = cor_texto

def main():
    """Função principal"""
    print("🚀 Criando apresentação PowerPoint: 'Como a LIA Revoluciona o Recrutamento Moderno'")
    print("📊 Processando conteúdo institucional...")

    try:
        arquivo = criar_apresentacao_lia()
        print(f"✅ Sucesso! Arquivo criado: {arquivo}")
        print("📥 Você pode baixar o arquivo .pptx agora")
        print("🎯 A apresentação contém 15 slides otimizados para conversão")
        return arquivo

    except Exception as e:
        print(f"❌ Erro ao criar apresentação: {str(e)}")
        return None

if __name__ == "__main__":
    main()
